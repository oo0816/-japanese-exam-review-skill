"""Parse PPTX files: extract text + images → JSON."""
import json
import sys
import os
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import io

OUTPUT_IMAGES = Path("output/images")
OUTPUT_IMAGES.mkdir(parents=True, exist_ok=True)


def extract_text_from_shape(shape):
    """Recursively extract all text from a shape."""
    texts = []
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if text:
                texts.append(text)
    if shape.has_table:
        table_data = []
        for row in shape.table.rows:
            row_texts = [cell.text.strip() for cell in row.cells]
            table_data.append(row_texts)
        texts.append({"type": "table", "data": table_data})
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            texts.extend(extract_text_from_shape(child))
    return texts


def extract_images_from_shape(shape, slide_num, img_counter):
    """Recursively extract images from a shape. Returns list of image info dicts."""
    images = []
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            image = shape.image
            ext = image.content_type.split("/")[-1]
            if ext == "jpeg":
                ext = "jpg"
            img_name = f"slide{slide_num}_img{img_counter[0]}.{ext}"
            img_path = OUTPUT_IMAGES / img_name

            # Save the image blob
            with open(img_path, "wb") as f:
                f.write(image.blob)

            # Convert non-PNG to PNG for better compatibility with fpdf2
            if ext != "png":
                png_name = f"slide{slide_num}_img{img_counter[0]}.png"
                png_path = OUTPUT_IMAGES / png_name
                try:
                    img = Image.open(io.BytesIO(image.blob))
                    img.save(png_path, "PNG")
                    img_name = png_name
                    # Remove original non-PNG file
                    if img_path != png_path:
                        img_path.unlink(missing_ok=True)
                except Exception:
                    pass  # keep original if conversion fails

            images.append({
                "file": img_name,
                "width_px": image.size[0] if image.size else 0,
                "height_px": image.size[1] if image.size else 0,
                "shape_name": shape.name
            })
            img_counter[0] += 1
        except Exception:
            pass
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            images.extend(extract_images_from_shape(child, slide_num, img_counter))
    return images


def parse_ppt(filepath):
    """Parse a single PPTX file and return structured data."""
    prs = Presentation(str(filepath))
    filename = os.path.basename(filepath)
    slides_data = []
    global_img_counter = [1]

    for slide_num, slide in enumerate(prs.slides, 1):
        texts = []
        tables = []
        images_data = []

        # First find the title
        title = ""
        if slide.shapes.title and slide.shapes.title.text.strip():
            title = slide.shapes.title.text.strip()

        for shape in slide.shapes:
            extracted = extract_text_from_shape(shape)
            for item in extracted:
                if isinstance(item, dict) and item.get("type") == "table":
                    tables.append(item["data"])
                elif isinstance(item, str):
                    # Don't duplicate the title
                    if item != title and item:
                        texts.append(item)

            images_data.extend(
                extract_images_from_shape(shape, slide_num, global_img_counter)
            )

        slides_data.append({
            "slide_num": slide_num,
            "title": title,
            "text_blocks": texts,
            "tables": tables,
            "images": images_data
        })

    return {"filename": filename, "slide_count": len(slides_data), "slides": slides_data}


def main():
    ppt_dir = Path("ppt")
    if not ppt_dir.exists():
        print(json.dumps({"error": "ppt/ directory not found"}, ensure_ascii=False))
        sys.exit(1)

    ppt_files = sorted(ppt_dir.glob("*.pptx"))
    if not ppt_files:
        print(json.dumps({"error": "No .pptx files found in ppt/"}, ensure_ascii=False))
        sys.exit(1)

    all_results = []
    for ppt_file in ppt_files:
        result = parse_ppt(ppt_file)
        all_results.append(result)

    output_path = Path("output/parsed_ppt.json")
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(all_results, f, ensure_ascii=False, indent=2)

    total_images = sum(
        sum(len(s.get("images", [])) for s in r["slides"])
        for r in all_results
    )
    print(json.dumps({
        "status": "ok",
        "files_parsed": len(all_results),
        "total_slides": sum(r["slide_count"] for r in all_results),
        "total_images_extracted": total_images,
        "output": str(output_path)
    }, ensure_ascii=False))


if __name__ == "__main__":
    main()
